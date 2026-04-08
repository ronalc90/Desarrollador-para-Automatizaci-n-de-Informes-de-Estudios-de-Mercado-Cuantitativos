"""Genera fixtures mínimos para los tests.

Crea programáticamente:
- ``sample_data.xlsx`` con dos hojas y dos tablas.
- ``sample_template.pptx`` con un slide que contiene un gráfico de
  columnas cuyo nombre es ``"Grafico satisfaccion"``.
- ``sample_mapping.yaml`` consistente con los dos archivos anteriores.

Se exportan las funciones ``ensure_fixtures`` y ``build_sample_xlsx``
para reutilizarlas desde los tests cuando sea necesario.
"""

from __future__ import annotations

from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


# ---------------------------------------------------------------------- #
# Datos canónicos que usan los tests                                     #
# ---------------------------------------------------------------------- #

SATISFACCION_HEADER = ["Dimension", "2022", "2023", "2024", "2025"]
SATISFACCION_ROWS = [
    ["Producto", 70, 72, 75, 78],
    ["Precio", 60, 62, 65, 66],
    ["Atencion", 80, 83, 85, 87],
    ["Entrega", 55, 58, 60, 64],
    ["Post-venta", 50, 53, 57, 60],
]

NPS_HEADER = ["Segmento", "Promotores", "Detractores"]
NPS_ROWS = [
    ["General", 55, 20],
    ["Jovenes", 60, 15],
    ["Adultos", 50, 25],
]

RECO_LEFT_HEADER = ["Canal", "Uso"]
RECO_LEFT_ROWS = [
    ["Web", 40],
    ["App", 55],
    ["Tienda", 30],
]

RECO_RIGHT_HEADER = ["Canal", "Score"]
RECO_RIGHT_ROWS = [
    ["Web", 7],
    ["App", 9],
    ["Tienda", 6],
]


def build_sample_xlsx(path: Path) -> Path:
    """Crea el Excel de prueba con varias hojas y tablas conocidas."""
    wb = Workbook()

    # Hoja 1: satisfacción.
    ws1 = wb.active
    ws1.title = "P1_satisfaccion"
    ws1.append(SATISFACCION_HEADER)
    for row in SATISFACCION_ROWS:
        ws1.append(row)

    # Hoja 2: NPS.
    ws2 = wb.create_sheet("P3_nps")
    ws2.append(NPS_HEADER)
    for row in NPS_ROWS:
        ws2.append(row)

    # Hoja 3: recomendación con dos tablas lado a lado (A:B y D:E).
    ws3 = wb.create_sheet("P4_recomendacion")
    for r, row in enumerate([RECO_LEFT_HEADER, *RECO_LEFT_ROWS], start=1):
        for c, value in enumerate(row, start=1):
            ws3.cell(row=r, column=c, value=value)
    for r, row in enumerate([RECO_RIGHT_HEADER, *RECO_RIGHT_ROWS], start=1):
        for c, value in enumerate(row, start=4):  # columnas D y E
            ws3.cell(row=r, column=c, value=value)

    # Hoja vacía para probar validaciones.
    wb.create_sheet("P99_vacia")

    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(path)
    return path


def build_sample_pptx(path: Path) -> Path:
    """Crea un template mínimo con varios gráficos.

    - Slide 1: portada vacía.
    - Slide 2: un gráfico con nombre conocido ("Grafico satisfaccion").
    - Slide 3: dos gráficos sin nombre único para probar matching por
      ``chart_index``.
    """
    prs = Presentation()

    # Slide 1: portada sin gráficos.
    portada_layout = prs.slide_layouts[0]
    slide_portada = prs.slides.add_slide(portada_layout)
    if slide_portada.shapes.title:
        slide_portada.shapes.title.text = "Estudio de ejemplo"

    # Slide 2: slide con un gráfico con nombre.
    blank_layout = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(blank_layout)

    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Serie inicial", (1, 2, 3))

    chart_shape = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.5),
        Inches(8),
        Inches(4.5),
        chart_data,
    )
    chart_shape.name = "Grafico satisfaccion"

    # Slide 3: dos gráficos para probar matching por chart_index.
    slide3 = prs.slides.add_slide(blank_layout)

    left_data = CategoryChartData()
    left_data.categories = ["X", "Y"]
    left_data.add_series("Serie inicial", (10, 20))
    slide3.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5),
        Inches(1.5),
        Inches(4.5),
        Inches(4.5),
        left_data,
    )

    right_data = CategoryChartData()
    right_data.categories = ["X", "Y"]
    right_data.add_series("Serie inicial", (5, 15))
    slide3.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(5.0),
        Inches(1.5),
        Inches(4.5),
        Inches(4.5),
        right_data,
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path


def build_sample_mapping(path: Path) -> Path:
    """Mapping mínimo compatible con el fixture base."""
    content = (
        "slides:\n"
        "  - slide_index: 2\n"
        "    charts:\n"
        "      - chart_name: \"Grafico satisfaccion\"\n"
        "        excel_sheet: \"P1_satisfaccion\"\n"
        "        data_range: \"A1:E6\"\n"
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


def build_sample_mapping_broken(path: Path) -> Path:
    """Mapping con errores intencionales para demos de validacion."""
    content = (
        "slides:\n"
        "  - slide_index: 99\n"
        "    charts:\n"
        "      - chart_name: \"Grafico que no existe\"\n"
        "        excel_sheet: \"P1_satisfaccion\"\n"
        "        data_range: \"A1:E6\"\n"
        "  - slide_index: 2\n"
        "    charts:\n"
        "      - chart_name: \"Grafico satisfaccion\"\n"
        "        excel_sheet: \"Hoja_inexistente\"\n"
        "        data_range: \"A1:E6\"\n"
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


def build_sample_mapping_multi(path: Path) -> Path:
    """Mapping extendido con matching por chart_index en el slide 3."""
    content = (
        "slides:\n"
        "  - slide_index: 2\n"
        "    charts:\n"
        "      - chart_name: \"Grafico satisfaccion\"\n"
        "        excel_sheet: \"P1_satisfaccion\"\n"
        "        data_range: \"A1:E6\"\n"
        "  - slide_index: 3\n"
        "    charts:\n"
        "      - chart_index: 0\n"
        "        excel_sheet: \"P4_recomendacion\"\n"
        "        data_range: \"A1:B4\"\n"
        "      - chart_index: 1\n"
        "        excel_sheet: \"P4_recomendacion\"\n"
        "        data_range: \"D1:E4\"\n"
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


# Versión del esquema de fixtures. Se incrementa cuando cambia la
# cantidad/forma de los fixtures generados para que ``ensure_fixtures``
# pueda detectar archivos viejos y regenerarlos.
FIXTURES_SCHEMA_VERSION = 4
_SCHEMA_MARKER = "fixtures.version"


# ---------------------------------------------------------------------- #
# Fixtures Etapa 3: procesamiento previo                                  #
# ---------------------------------------------------------------------- #

RESPONSES_ROWS = [
    # respondent_id, segmento, edad, satisfaccion (1-5), nps (0-10), gasto
    (1, "Jovenes", 22, 5, 10, 120),
    (2, "Jovenes", 25, 4, 8, 95),
    (3, "Jovenes", 28, 3, 6, 80),
    (4, "Jovenes", 31, 5, 9, 150),
    (5, "Adultos", 35, 4, 7, 200),
    (6, "Adultos", 40, 3, 5, 180),
    (7, "Adultos", 45, 2, 3, 90),
    (8, "Adultos", 50, 5, 10, 250),
    (9, "Adultos", 55, 4, 8, 210),
    (10, "Mayores", 60, 3, 6, 170),
    (11, "Mayores", 65, 2, 4, 100),
    (12, "Mayores", 70, 4, 7, 140),
    (13, "Mayores", 72, 5, 9, 160),
    (14, "Jovenes", 24, 4, 8, 130),
    (15, "Adultos", 42, 5, 10, 220),
]


def build_sample_responses_csv(path: Path) -> Path:
    """CSV con datos crudos de respuestas de un estudio simulado."""
    path.parent.mkdir(parents=True, exist_ok=True)
    header = "respondent_id,segmento,edad,satisfaccion,nps,gasto\n"
    lines = [
        ",".join(str(v) for v in row) for row in RESPONSES_ROWS
    ]
    path.write_text(header + "\n".join(lines) + "\n", encoding="utf-8")
    return path


def build_sample_tab_plan(path: Path) -> Path:
    """Tab Plan en YAML compatible con el CSV de respuestas."""
    content = (
        "version: 1\n"
        "crosses:\n"
        "  - name: \"Satisfaccion por segmento\"\n"
        "    rows: [satisfaccion]\n"
        "    columns: [segmento]\n"
        "  - name: \"NPS promedio por segmento\"\n"
        "    rows: [segmento]\n"
        "    values: nps\n"
        "    aggregate: mean\n"
        "  - name: \"Gasto total por segmento\"\n"
        "    rows: [segmento]\n"
        "    values: gasto\n"
        "    aggregate: sum\n"
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


def build_sample_tab_plan_xlsx(path: Path) -> Path:
    """Tab Plan embebido en la hoja 'TabPlan' de un xlsx."""
    wb = Workbook()
    ws = wb.active
    ws.title = "TabPlan"
    ws.append(["name", "rows", "columns", "values", "aggregate", "percentage"])
    ws.append([
        "Satisfaccion por segmento",
        "satisfaccion",
        "segmento",
        None,
        "count",
        "none",
    ])
    ws.append([
        "NPS promedio por segmento",
        "segmento",
        None,
        "nps",
        "mean",
        "none",
    ])
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(path)
    return path


def _read_schema_version(fixtures_dir: Path) -> int:
    marker = fixtures_dir / _SCHEMA_MARKER
    if not marker.exists():
        return 0
    try:
        return int(marker.read_text(encoding="utf-8").strip())
    except ValueError:
        return 0


def _write_schema_version(fixtures_dir: Path) -> None:
    marker = fixtures_dir / _SCHEMA_MARKER
    marker.write_text(str(FIXTURES_SCHEMA_VERSION), encoding="utf-8")


def ensure_fixtures(fixtures_dir: Path) -> None:
    """Se asegura de que existan los fixtures y sean de la version actual.

    Si la marca de versión no existe o es vieja, regenera todos los
    archivos para que coincidan con el schema de los tests actuales.
    """
    fixtures_dir.mkdir(parents=True, exist_ok=True)
    xlsx = fixtures_dir / "sample_data.xlsx"
    pptx = fixtures_dir / "sample_template.pptx"
    mapping = fixtures_dir / "sample_mapping.yaml"
    mapping_multi = fixtures_dir / "sample_mapping_multi.yaml"
    mapping_broken = fixtures_dir / "sample_mapping_broken.yaml"

    # Etapa 3.
    responses_csv = fixtures_dir / "sample_responses.csv"
    tab_plan_yaml = fixtures_dir / "sample_tab_plan.yaml"
    tab_plan_xlsx = fixtures_dir / "sample_tab_plan.xlsx"

    current_version = _read_schema_version(fixtures_dir)
    stale = current_version < FIXTURES_SCHEMA_VERSION

    if stale or not xlsx.exists():
        build_sample_xlsx(xlsx)
    if stale or not pptx.exists():
        build_sample_pptx(pptx)
    if stale or not mapping.exists():
        build_sample_mapping(mapping)
    if stale or not mapping_multi.exists():
        build_sample_mapping_multi(mapping_multi)
    if stale or not mapping_broken.exists():
        build_sample_mapping_broken(mapping_broken)
    if stale or not responses_csv.exists():
        build_sample_responses_csv(responses_csv)
    if stale or not tab_plan_yaml.exists():
        build_sample_tab_plan(tab_plan_yaml)
    if stale or not tab_plan_xlsx.exists():
        build_sample_tab_plan_xlsx(tab_plan_xlsx)

    _write_schema_version(fixtures_dir)


if __name__ == "__main__":
    target = Path(__file__).parent
    ensure_fixtures(target)
    print(f"Fixtures generados en {target}")
