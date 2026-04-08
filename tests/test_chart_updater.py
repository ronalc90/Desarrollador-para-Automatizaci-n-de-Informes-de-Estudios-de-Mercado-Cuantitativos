"""Tests unitarios para ``engine.chart_updater``."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
import pytest
from pptx import Presentation

from engine.chart_updater import (
    ChartDataShapeError,
    find_chart_in_slide,
    update_chart_with_dataframe,
)


def _get_first_chart(path: Path):
    prs = Presentation(str(path))
    slide = prs.slides[1]  # el fixture tiene el grafico en el slide 2
    chart = find_chart_in_slide(slide, chart_name="Grafico satisfaccion")
    return prs, slide, chart


def test_find_chart_by_name(sample_pptx: Path) -> None:
    _, _, chart = _get_first_chart(sample_pptx)
    assert chart is not None


def test_find_chart_by_index(sample_pptx: Path) -> None:
    prs = Presentation(str(sample_pptx))
    slide = prs.slides[1]
    chart = find_chart_in_slide(slide, chart_index=0)
    assert chart is not None


def test_find_chart_not_found(sample_pptx: Path) -> None:
    prs = Presentation(str(sample_pptx))
    slide = prs.slides[0]  # portada sin graficos
    assert find_chart_in_slide(slide, chart_name="X") is None


def test_update_chart_with_dataframe(
    tmp_path: Path, sample_pptx: Path
) -> None:
    df = pd.DataFrame(
        {
            "Dimension": ["A", "B", "C"],
            "2024": [10, 20, 30],
            "2025": [15, 25, 35],
        }
    )

    prs, _, chart = _get_first_chart(sample_pptx)
    update_chart_with_dataframe(chart, df)

    out = tmp_path / "updated.pptx"
    prs.save(str(out))

    # Reabrir y verificar que las categorias y valores quedaron.
    reopened = Presentation(str(out))
    chart2 = find_chart_in_slide(
        reopened.slides[1], chart_name="Grafico satisfaccion"
    )
    assert chart2 is not None
    categories = [c for c in chart2.plots[0].categories]
    assert categories == ["A", "B", "C"]
    series_names = [s.name for s in chart2.series]
    assert series_names == ["2024", "2025"]
    # Los valores de la primera serie deben coincidir con el DF.
    assert list(chart2.series[0].values) == [10, 20, 30]
    assert list(chart2.series[1].values) == [15, 25, 35]


def test_update_chart_with_empty_df(sample_pptx: Path) -> None:
    _, _, chart = _get_first_chart(sample_pptx)
    with pytest.raises(ChartDataShapeError):
        update_chart_with_dataframe(chart, pd.DataFrame())


def test_update_chart_with_one_column_df(sample_pptx: Path) -> None:
    _, _, chart = _get_first_chart(sample_pptx)
    with pytest.raises(ChartDataShapeError):
        update_chart_with_dataframe(chart, pd.DataFrame({"a": [1, 2]}))
