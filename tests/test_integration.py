"""Test de integración end-to-end del pipeline completo."""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation

from engine.batch_processor import process_batch
from engine.chart_updater import find_chart_in_slide
from engine.ppt_builder import build_presentation
from engine.validator import validate_all


def test_pipeline_end_to_end(
    tmp_path: Path,
    sample_pptx: Path,
    sample_xlsx: Path,
    sample_mapping: Path,
) -> None:
    """Ejecuta el pipeline y verifica el archivo de salida."""
    output_dir = tmp_path / "out"

    # Primero validación.
    result = validate_all(sample_pptx, sample_xlsx, sample_mapping)
    assert result.ok, result.as_report()

    # Build.
    build_result = build_presentation(
        template_path=sample_pptx,
        excel_path=sample_xlsx,
        mapping=sample_mapping,
        output_dir=output_dir,
    )
    assert build_result.ok
    assert build_result.charts_updated == 1
    assert build_result.output_path.exists()

    # Verificar que el XML del grafico tiene los nuevos valores.
    prs = Presentation(str(build_result.output_path))
    chart = find_chart_in_slide(
        prs.slides[1], chart_name="Grafico satisfaccion"
    )
    assert chart is not None
    categories = [c for c in chart.plots[0].categories]
    assert categories == [
        "Producto",
        "Precio",
        "Atencion",
        "Entrega",
        "Post-venta",
    ]
    series_names = [s.name for s in chart.series]
    assert series_names == ["2022", "2023", "2024", "2025"]

    # Verificar que el Excel embebido tambien fue actualizado.
    with zipfile.ZipFile(str(build_result.output_path), "r") as zf:
        embedded_names = [
            n for n in zf.namelist() if n.startswith("ppt/embeddings/")
        ]
        assert embedded_names, "No hay archivo embebido en el pptx."
        with zf.open(embedded_names[0]) as embedded:
            wb = load_workbook(io.BytesIO(embedded.read()))
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    # Primera fila: header.
    assert rows[0][0] == "Dimension"
    assert rows[0][-1] == "2025"
    # Primera fila de datos.
    assert rows[1][0] == "Producto"
    assert rows[1][-1] == 78


def test_batch_pipeline(
    tmp_path: Path,
    sample_pptx: Path,
    sample_xlsx: Path,
    sample_mapping: Path,
) -> None:
    """El batch debe procesar varios archivos con el mismo template."""
    data_folder = tmp_path / "data"
    data_folder.mkdir()
    # Copiamos el mismo xlsx varias veces con distintos nombres.
    for name in ("estudio_chile.xlsx", "estudio_peru.xlsx", "estudio_ar.xlsx"):
        (data_folder / name).write_bytes(sample_xlsx.read_bytes())

    output_dir = tmp_path / "out"
    result = process_batch(
        template_path=sample_pptx,
        data_folder=data_folder,
        mapping=sample_mapping,
        output_dir=output_dir,
    )
    assert len(result.items) == 3
    assert len(result.failed) == 0
    assert len(result.successful) == 3
    for item in result.successful:
        assert item.build_result is not None
        assert item.build_result.output_path.exists()


def test_pipeline_with_chart_index_matching(
    tmp_path: Path,
    sample_pptx: Path,
    sample_xlsx: Path,
    sample_mapping_multi: Path,
) -> None:
    """Verifica matching por chart_index en un slide con varios graficos."""
    result = validate_all(sample_pptx, sample_xlsx, sample_mapping_multi)
    assert result.ok, result.as_report()

    build_result = build_presentation(
        template_path=sample_pptx,
        excel_path=sample_xlsx,
        mapping=sample_mapping_multi,
        output_dir=tmp_path / "out",
    )
    assert build_result.ok
    # 1 grafico en slide 2 + 2 graficos en slide 3 = 3 graficos actualizados.
    assert build_result.charts_updated == 3

    prs = Presentation(str(build_result.output_path))

    # Slide 3 contiene los dos graficos por chart_index.
    slide3 = prs.slides[2]
    charts = [s.chart for s in slide3.shapes if s.has_chart]
    assert len(charts) == 2

    # chart_index 0 -> tabla izquierda (Canal, Uso)
    left_categories = [c for c in charts[0].plots[0].categories]
    assert left_categories == ["Web", "App", "Tienda"]
    assert charts[0].series[0].name == "Uso"
    assert list(charts[0].series[0].values) == [40, 55, 30]

    # chart_index 1 -> tabla derecha (Canal, Score)
    right_categories = [c for c in charts[1].plots[0].categories]
    assert right_categories == ["Web", "App", "Tienda"]
    assert charts[1].series[0].name == "Score"
    assert list(charts[1].series[0].values) == [7, 9, 6]
