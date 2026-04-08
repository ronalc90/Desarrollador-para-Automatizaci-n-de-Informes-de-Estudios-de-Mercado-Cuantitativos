"""Construcción del PPT final a partir de un template y un Excel.

Este es el orquestador principal del motor. Toma un template .pptx,
un Excel de datos y un mapping YAML, y produce un archivo .pptx en
el directorio de salida con todos los gráficos actualizados.

Diseño
------
- Carga el template con python-pptx.
- Carga el mapping con :func:`engine.validator.load_mapping`.
- Abre el Excel con :class:`engine.excel_reader.ExcelReader`.
- Para cada slide del mapping, busca los gráficos por nombre/indice,
  carga los datos con ``reader.get_table`` y actualiza cada gráfico
  con :func:`engine.chart_updater.update_chart_with_dataframe`.
- Los slides sin gráficos (portadas, separadores) se preservan.
- Si un gráfico individual falla, se loguea y se continúa con los
  siguientes. El archivo final siempre se guarda, incluso si hubo
  gráficos que no pudieron actualizarse.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional, Union

from pptx import Presentation

from engine.chart_updater import (
    ChartUpdaterError,
    find_chart_in_slide,
    update_chart_with_dataframe,
)
from engine.excel_reader import ExcelReader, ExcelReaderError
from engine.validator import Mapping, load_mapping

logger = logging.getLogger(__name__)


class PPTBuilderError(Exception):
    """Error base del módulo ppt_builder."""


@dataclass
class BuildResult:
    """Resultado de un ``build_presentation``."""

    output_path: Path
    charts_updated: int = 0
    charts_failed: int = 0
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return self.charts_failed == 0 and not self.errors

    def summary(self) -> str:
        parts = [
            f"Archivo: {self.output_path}",
            f"Gráficos actualizados: {self.charts_updated}",
            f"Gráficos fallidos: {self.charts_failed}",
        ]
        if self.warnings:
            parts.append("Warnings:")
            parts.extend(f"  - {w}" for w in self.warnings)
        if self.errors:
            parts.append("Errores:")
            parts.extend(f"  - {e}" for e in self.errors)
        return "\n".join(parts)


def _slugify(name: str) -> str:
    keep = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789-_"
    cleaned = "".join(c if c in keep else "_" for c in name)
    return cleaned.strip("_") or "output"


def _resolve_output_path(
    output_dir: Path,
    excel_path: Path,
    output_name: Optional[str],
) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    if output_name:
        if not output_name.lower().endswith(".pptx"):
            output_name = f"{output_name}.pptx"
        return output_dir / output_name

    slug = _slugify(excel_path.stem)
    return output_dir / f"output_{slug}.pptx"


def build_presentation(
    template_path: Union[str, Path],
    excel_path: Union[str, Path],
    mapping: Union[str, Path, Mapping],
    output_dir: Union[str, Path],
    output_name: Optional[str] = None,
) -> BuildResult:
    """Genera un PPT final aplicando ``excel_path`` sobre ``template_path``.

    Parameters
    ----------
    template_path:
        Ruta al template .pptx.
    excel_path:
        Ruta al Excel .xlsx con los datos.
    mapping:
        Ruta al YAML del mapping o un objeto :class:`Mapping` ya cargado.
    output_dir:
        Directorio donde se escribirá el archivo final (se crea si no
        existe).
    output_name:
        Nombre de archivo de salida. Si es ``None``, se deriva del nombre
        del Excel como ``output_{nombre}.pptx``.

    Returns
    -------
    BuildResult
        Resumen del proceso con contadores y errores.
    """
    template_path = Path(template_path)
    excel_path = Path(excel_path)
    output_dir = Path(output_dir)

    if not template_path.exists():
        raise PPTBuilderError(f"Template no encontrado: {template_path}")
    if not excel_path.exists():
        raise PPTBuilderError(f"Excel no encontrado: {excel_path}")

    if isinstance(mapping, (str, Path)):
        mapping_obj = load_mapping(mapping)
    elif isinstance(mapping, Mapping):
        mapping_obj = mapping
    else:
        raise PPTBuilderError(
            "El parámetro 'mapping' debe ser una ruta o un Mapping."
        )

    output_path = _resolve_output_path(output_dir, excel_path, output_name)
    result = BuildResult(output_path=output_path)

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        raise PPTBuilderError(
            f"No se pudo abrir el template {template_path}: {exc}"
        ) from exc

    try:
        reader = ExcelReader(excel_path)
    except ExcelReaderError as exc:
        raise PPTBuilderError(
            f"No se pudo abrir el Excel {excel_path}: {exc}"
        ) from exc

    try:
        total_slides = len(prs.slides)

        for slide_map in mapping_obj.slides:
            idx0 = slide_map.slide_index - 1
            if idx0 < 0 or idx0 >= total_slides:
                msg = (
                    f"slide_index {slide_map.slide_index} fuera de rango "
                    f"(el template tiene {total_slides} slides). "
                    "Se omite."
                )
                logger.warning(msg)
                result.warnings.append(msg)
                continue

            slide = prs.slides[idx0]

            for chart_map in slide_map.charts:
                chart = find_chart_in_slide(
                    slide,
                    chart_name=chart_map.chart_name,
                    chart_index=chart_map.chart_index,
                )
                if chart is None:
                    msg = (
                        f"slide {slide_map.slide_index}: no se encontró el "
                        f"gráfico (name={chart_map.chart_name!r}, "
                        f"index={chart_map.chart_index}). Se omite."
                    )
                    logger.warning(msg)
                    result.warnings.append(msg)
                    result.charts_failed += 1
                    continue

                try:
                    df = reader.get_table(
                        chart_map.excel_sheet,
                        chart_map.data_range,
                    )
                except ExcelReaderError as exc:
                    msg = (
                        f"slide {slide_map.slide_index} / "
                        f"{chart_map.excel_sheet}!{chart_map.data_range}: "
                        f"no se pudo leer el Excel ({exc})."
                    )
                    logger.error(msg)
                    result.errors.append(msg)
                    result.charts_failed += 1
                    continue

                try:
                    update_chart_with_dataframe(chart, df)
                except ChartUpdaterError as exc:
                    msg = (
                        f"slide {slide_map.slide_index}: no se pudo "
                        f"actualizar el gráfico "
                        f"{chart_map.chart_name or chart_map.chart_index}: "
                        f"{exc}"
                    )
                    logger.error(msg)
                    result.errors.append(msg)
                    result.charts_failed += 1
                    continue

                result.charts_updated += 1
                logger.info(
                    "Actualizado gráfico en slide %s con %s!%s",
                    slide_map.slide_index,
                    chart_map.excel_sheet,
                    chart_map.data_range,
                )
    finally:
        reader.close()

    try:
        prs.save(str(output_path))
    except Exception as exc:
        raise PPTBuilderError(
            f"No se pudo guardar el PPT en {output_path}: {exc}"
        ) from exc

    return result


__all__ = ["BuildResult", "PPTBuilderError", "build_presentation"]
