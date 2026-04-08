"""Inspección de templates .pptx.

Módulo de diagnóstico para ayudar al usuario a armar el ``mapping.yaml``
sin adivinar. Dada una ruta a un ``.pptx`` devuelve (y sabe imprimir)
la lista de slides, la cantidad de gráficos por slide y los nombres de
shape de cada gráfico.

Uso típico desde el CLI::

    python main.py inspect --template template.pptx

O programáticamente::

    report = inspect_template("template.pptx")
    print(report.as_text())
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Union

from pptx import Presentation


@dataclass
class ChartInfo:
    """Información mínima de un gráfico dentro de un slide."""

    chart_index: int
    shape_name: str
    chart_type: str


@dataclass
class SlideInfo:
    """Información de un slide del template."""

    slide_index: int  # 1-based, listo para el mapping
    layout_name: str
    charts: list[ChartInfo] = field(default_factory=list)

    @property
    def has_charts(self) -> bool:
        return bool(self.charts)


@dataclass
class TemplateReport:
    """Reporte completo de inspección de un template."""

    template_path: Path
    slides: list[SlideInfo] = field(default_factory=list)

    @property
    def total_slides(self) -> int:
        return len(self.slides)

    @property
    def total_charts(self) -> int:
        return sum(len(s.charts) for s in self.slides)

    def as_text(self) -> str:
        """Serializa el reporte a texto plano amigable para consola."""
        lines = [
            f"Template: {self.template_path}",
            f"Total slides: {self.total_slides}",
            f"Total graficos: {self.total_charts}",
            "",
        ]
        for slide in self.slides:
            marker = "*" if slide.has_charts else " "
            lines.append(
                f"{marker} Slide {slide.slide_index} "
                f"[{slide.layout_name}] - {len(slide.charts)} grafico(s)"
            )
            for chart in slide.charts:
                lines.append(
                    f"    - chart_index={chart.chart_index}, "
                    f"chart_name={chart.shape_name!r}, "
                    f"tipo={chart.chart_type}"
                )
        return "\n".join(lines)

    def as_mapping_stub(self) -> str:
        """Genera un stub de mapping.yaml listo para editar.

        Solo incluye los slides con gráficos. Deja los campos
        ``excel_sheet`` y ``data_range`` como placeholders para que el
        usuario los complete.
        """
        lines = ["slides:"]
        any_chart = False
        for slide in self.slides:
            if not slide.has_charts:
                continue
            any_chart = True
            lines.append(f"  - slide_index: {slide.slide_index}")
            lines.append("    charts:")
            for chart in slide.charts:
                lines.append(f"      - chart_name: {chart.shape_name!r}")
                lines.append("        excel_sheet: TODO")
                lines.append("        data_range: TODO")
        if not any_chart:
            lines.append("  []  # el template no tiene graficos")
        return "\n".join(lines)


def _chart_type_name(chart) -> str:
    """Devuelve un nombre humano para el tipo de gráfico."""
    try:
        return chart.chart_type.name if chart.chart_type is not None else "?"
    except Exception:
        return "?"


def inspect_template(template_path: Union[str, Path]) -> TemplateReport:
    """Devuelve un :class:`TemplateReport` del template indicado.

    Parameters
    ----------
    template_path:
        Ruta a un archivo .pptx existente.

    Raises
    ------
    FileNotFoundError
        Si el template no existe.
    ValueError
        Si el archivo no puede abrirse como .pptx.
    """
    path = Path(template_path)
    if not path.exists():
        raise FileNotFoundError(f"Template no encontrado: {path}")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ValueError(f"No se pudo abrir {path}: {exc}") from exc

    report = TemplateReport(template_path=path)

    for idx, slide in enumerate(prs.slides, start=1):
        layout_name = getattr(slide.slide_layout, "name", "sin_layout") or "sin_layout"
        slide_info = SlideInfo(slide_index=idx, layout_name=layout_name)

        chart_idx = 0
        for shape in slide.shapes:
            if shape.has_chart:
                slide_info.charts.append(
                    ChartInfo(
                        chart_index=chart_idx,
                        shape_name=shape.name,
                        chart_type=_chart_type_name(shape.chart),
                    )
                )
                chart_idx += 1

        report.slides.append(slide_info)

    return report


__all__ = [
    "ChartInfo",
    "SlideInfo",
    "TemplateReport",
    "inspect_template",
]
