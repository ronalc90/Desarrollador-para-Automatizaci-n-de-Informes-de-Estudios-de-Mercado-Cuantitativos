"""Validaciones previas al procesamiento.

El validador se ejecuta ``fail-fast``: antes de tocar cualquier archivo
de salida revisa que el mapping, el Excel de datos y el template .pptx
sean compatibles entre sí. Devuelve un ``ValidationResult`` con listas
de errores (bloqueantes) y warnings (informativos).

Reglas principales
------------------

1. El YAML de mapping tiene la estructura esperada
   (``slides -> charts -> excel_sheet + data_range``).
2. Cada ``excel_sheet`` referenciada existe en el Excel.
3. Cada ``data_range`` referenciado no está vacío.
4. Cada ``slide_index`` referenciado existe en el template.
5. Los gráficos referenciados por nombre existen en ese slide (warning
   si no se encuentra el nombre, porque se permite fallback por índice).
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable, Optional, Union

import yaml
from pptx import Presentation

from engine.excel_reader import (
    ExcelReader,
    ExcelReaderError,
    RangeNotFoundError,
    SheetNotFoundError,
)


# ---------------------------------------------------------------------- #
# Modelos                                                                #
# ---------------------------------------------------------------------- #


@dataclass
class ChartMapping:
    """Una entrada ``chart`` dentro de un slide del mapping."""

    chart_name: Optional[str]
    excel_sheet: str
    data_range: str
    chart_index: Optional[int] = None


@dataclass
class SlideMapping:
    """Una entrada ``slide`` del mapping."""

    slide_index: int
    charts: list[ChartMapping]


@dataclass
class Mapping:
    """Mapping completo cargado desde YAML."""

    slides: list[SlideMapping]

    @classmethod
    def from_dict(cls, data: dict[str, Any]) -> "Mapping":
        if not isinstance(data, dict):
            raise ValueError("El mapping debe ser un diccionario YAML.")
        raw_slides = data.get("slides")
        if not isinstance(raw_slides, list) or not raw_slides:
            raise ValueError(
                "El mapping debe tener una clave 'slides' con al menos una entrada."
            )

        slides: list[SlideMapping] = []
        for idx, raw in enumerate(raw_slides):
            if not isinstance(raw, dict):
                raise ValueError(f"slides[{idx}] debe ser un diccionario.")
            try:
                slide_index = int(raw["slide_index"])
            except (KeyError, TypeError, ValueError) as exc:
                raise ValueError(
                    f"slides[{idx}] debe tener un 'slide_index' entero."
                ) from exc

            raw_charts = raw.get("charts")
            if not isinstance(raw_charts, list) or not raw_charts:
                raise ValueError(
                    f"slides[{idx}] debe tener una lista 'charts' no vacía."
                )

            charts: list[ChartMapping] = []
            for c_idx, rc in enumerate(raw_charts):
                if not isinstance(rc, dict):
                    raise ValueError(
                        f"slides[{idx}].charts[{c_idx}] debe ser un diccionario."
                    )
                try:
                    sheet = str(rc["excel_sheet"])
                    data_range = str(rc["data_range"])
                except KeyError as exc:
                    raise ValueError(
                        f"slides[{idx}].charts[{c_idx}] debe tener "
                        "'excel_sheet' y 'data_range'."
                    ) from exc

                chart_name = rc.get("chart_name")
                chart_index = rc.get("chart_index")
                if chart_index is not None:
                    chart_index = int(chart_index)
                if chart_name is not None:
                    chart_name = str(chart_name)

                charts.append(
                    ChartMapping(
                        chart_name=chart_name,
                        excel_sheet=sheet,
                        data_range=data_range,
                        chart_index=chart_index,
                    )
                )
            slides.append(SlideMapping(slide_index=slide_index, charts=charts))

        return cls(slides=slides)


@dataclass
class ValidationResult:
    """Resultado de ejecutar el validador."""

    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return not self.errors

    def extend(self, other: "ValidationResult") -> None:
        self.errors.extend(other.errors)
        self.warnings.extend(other.warnings)

    def as_report(self) -> str:
        parts: list[str] = []
        if self.errors:
            parts.append("ERRORES:")
            parts.extend(f"  - {e}" for e in self.errors)
        if self.warnings:
            parts.append("WARNINGS:")
            parts.extend(f"  - {w}" for w in self.warnings)
        if not parts:
            parts.append("Sin errores ni warnings.")
        return "\n".join(parts)


# ---------------------------------------------------------------------- #
# API pública                                                            #
# ---------------------------------------------------------------------- #


def load_mapping(path: Union[str, Path]) -> Mapping:
    """Carga y valida el shape básico del mapping desde un archivo YAML."""
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"Mapping no encontrado: {path}")
    with path.open("r", encoding="utf-8") as fh:
        data = yaml.safe_load(fh)
    return Mapping.from_dict(data)


def validate_excel_against_mapping(
    excel_path: Union[str, Path],
    mapping: Mapping,
) -> ValidationResult:
    """Valida que el Excel tenga las hojas y rangos referenciados."""
    result = ValidationResult()
    try:
        reader = ExcelReader(excel_path)
    except ExcelReaderError as exc:
        result.errors.append(f"No se pudo abrir el Excel: {exc}")
        return result

    try:
        available_sheets = set(reader.sheet_names)
        # Comprobar cada (sheet, range) del mapping.
        for slide in mapping.slides:
            for chart in slide.charts:
                if chart.excel_sheet not in available_sheets:
                    result.errors.append(
                        f"slide {slide.slide_index}: hoja "
                        f"'{chart.excel_sheet}' no existe en el Excel."
                    )
                    continue
                try:
                    df = reader.get_table(chart.excel_sheet, chart.data_range)
                except SheetNotFoundError as exc:
                    result.errors.append(str(exc))
                except RangeNotFoundError as exc:
                    result.errors.append(
                        f"slide {slide.slide_index} / "
                        f"{chart.excel_sheet}!{chart.data_range}: {exc}"
                    )
                except ExcelReaderError as exc:
                    result.errors.append(
                        f"slide {slide.slide_index} / "
                        f"{chart.excel_sheet}!{chart.data_range}: {exc}"
                    )
                else:
                    if df.empty:
                        result.warnings.append(
                            f"slide {slide.slide_index} / "
                            f"{chart.excel_sheet}!{chart.data_range}: "
                            "la tabla no tiene filas de datos."
                        )
    finally:
        reader.close()

    return result


def validate_template_against_mapping(
    template_path: Union[str, Path],
    mapping: Mapping,
) -> ValidationResult:
    """Valida que el template .pptx contenga los slides/charts del mapping."""
    result = ValidationResult()
    template_path = Path(template_path)
    if not template_path.exists():
        result.errors.append(f"Template no encontrado: {template_path}")
        return result

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        result.errors.append(f"No se pudo abrir el template: {exc}")
        return result

    total_slides = len(prs.slides)

    for slide_map in mapping.slides:
        # slide_index del mapping es 1-based para ser amigable al usuario.
        idx0 = slide_map.slide_index - 1
        if idx0 < 0 or idx0 >= total_slides:
            result.errors.append(
                f"El slide_index {slide_map.slide_index} no existe "
                f"(el template tiene {total_slides} slides)."
            )
            continue

        slide = prs.slides[idx0]
        charts_in_slide = [
            shape for shape in slide.shapes if shape.has_chart
        ]
        if not charts_in_slide:
            result.errors.append(
                f"El slide {slide_map.slide_index} no tiene gráficos, "
                "pero el mapping lo referencia."
            )
            continue

        names_in_slide = {shape.name for shape in charts_in_slide}
        for chart_map in slide_map.charts:
            if chart_map.chart_name and chart_map.chart_name not in names_in_slide:
                result.warnings.append(
                    f"slide {slide_map.slide_index}: no se encontró un "
                    f"gráfico con nombre '{chart_map.chart_name}'. "
                    "Se usará match por índice si está disponible."
                )
            if (
                chart_map.chart_index is not None
                and chart_map.chart_index >= len(charts_in_slide)
            ):
                result.errors.append(
                    f"slide {slide_map.slide_index}: chart_index "
                    f"{chart_map.chart_index} fuera de rango "
                    f"(el slide tiene {len(charts_in_slide)} gráficos)."
                )

    return result


def validate_all(
    template_path: Union[str, Path],
    excel_path: Union[str, Path],
    mapping_path: Union[str, Path],
) -> ValidationResult:
    """Ejecuta todas las validaciones y agrega los resultados."""
    result = ValidationResult()
    try:
        mapping = load_mapping(mapping_path)
    except (FileNotFoundError, ValueError, yaml.YAMLError) as exc:
        result.errors.append(f"Mapping inválido: {exc}")
        return result

    result.extend(validate_excel_against_mapping(excel_path, mapping))
    result.extend(validate_template_against_mapping(template_path, mapping))
    return result


__all__ = [
    "ChartMapping",
    "SlideMapping",
    "Mapping",
    "ValidationResult",
    "load_mapping",
    "validate_excel_against_mapping",
    "validate_template_against_mapping",
    "validate_all",
]
